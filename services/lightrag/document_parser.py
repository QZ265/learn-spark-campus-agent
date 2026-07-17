import json
import re
import zipfile
from io import BytesIO
from xml.etree import ElementTree
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import List, Optional


SOURCE_MARKER = "PYLS_SOURCE_META"


@dataclass
class ParsedSection:
    chapter: str
    page_index: Optional[int]
    content: str
    source_kind: str


@dataclass
class ParsedDocument:
    title: str
    parser: str
    parser_version: str
    sections: List[ParsedSection]

    @property
    def text(self) -> str:
        return "\n\n".join(section.content for section in self.sections)

    def as_dict(self) -> dict:
        return {
            "title": self.title,
            "parser": self.parser,
            "parser_version": self.parser_version,
            "sections": [asdict(section) for section in self.sections],
        }


def _nonempty(value: str) -> str:
    return re.sub(r"\n{3,}", "\n\n", (value or "").strip())


def _parse_pdf(path: Path) -> ParsedDocument:
    from pypdf import PdfReader, __version__ as parser_version

    reader = PdfReader(str(path))
    sections = []
    for index, page in enumerate(reader.pages, start=1):
        text = _nonempty(page.extract_text() or "")
        source_kind = "page"
        if len(text) < 25:
            text = _ocr_pdf_page(path, index - 1)
            source_kind = "ocr_page"
        if text:
            first_line = text.splitlines()[0].strip()[:100]
            sections.append(ParsedSection(first_line or f"第 {index} 页", index, text, source_kind))
    if not sections:
        raise ValueError("PDF 没有可提取文本层，需要 OCR 后再导入")
    return ParsedDocument(path.stem, "pypdf", parser_version, sections)


def _ocr_pdf_page(path: Path, zero_based_page: int) -> str:
    try:
        import fitz
        from PIL import Image
        from ocrmac import ocrmac
    except ImportError as exc:
        raise ValueError("PDF 没有文本层且 OCR 依赖未安装") from exc
    try:
        with fitz.open(str(path)) as document:
            page = document.load_page(zero_based_page)
            pixmap = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
            image = Image.open(BytesIO(pixmap.tobytes("png")))
        results = ocrmac.OCR(
            image,
            recognition_level="accurate",
            language_preference=["zh-Hans", "en-US"],
            confidence_threshold=0.35,
        ).recognize()
        return _nonempty("\n".join(str(item[0]) for item in results if item and str(item[0]).strip()))
    except Exception as exc:
        raise ValueError(f"第 {zero_based_page + 1} 页 OCR 失败：{exc}") from exc


def _parse_docx(path: Path) -> ParsedDocument:
    import docx

    try:
        document = docx.Document(str(path))
    except (KeyError, ValueError, zipfile.BadZipFile):
        return _parse_docx_xml_fallback(path, getattr(docx, "__version__", "unknown"))
    sections: List[ParsedSection] = []
    current_title = path.stem
    current_lines: List[str] = []

    def flush() -> None:
        nonlocal current_lines
        text = _nonempty("\n".join(current_lines))
        if text:
            sections.append(ParsedSection(current_title, len(sections) + 1, text, "section"))
        current_lines = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style = (paragraph.style.name if paragraph.style else "").lower()
        if style.startswith("heading") or style.startswith("标题"):
            flush()
            current_title = text[:120]
        else:
            current_lines.append(text)
    for table in document.tables:
        rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
        current_lines.extend(row for row in rows if row.strip(" |"))
    flush()
    if not sections:
        raise ValueError("DOCX 中没有可提取文本")
    return ParsedDocument(path.stem, "python-docx", getattr(docx, "__version__", "unknown"), sections)


def _parse_docx_xml_fallback(path: Path, version: str) -> ParsedDocument:
    """Recover text from DOCX files whose optional media relationships are broken."""
    with zipfile.ZipFile(path) as archive:
        raw = archive.read("word/document.xml")
    root = ElementTree.fromstring(raw)
    namespace = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
    paragraphs = []
    for paragraph in root.iter(namespace + "p"):
        text = "".join(node.text or "" for node in paragraph.iter(namespace + "t")).strip()
        if text:
            paragraphs.append(text)
    if not paragraphs:
        raise ValueError("DOCX 正文 XML 中没有可提取文本")
    sections = []
    for offset in range(0, len(paragraphs), 35):
        block = paragraphs[offset : offset + 35]
        sections.append(ParsedSection(block[0][:100], len(sections) + 1, "\n".join(block), "section"))
    return ParsedDocument(path.stem, "python-docx-xml-fallback", version, sections)


def _parse_pptx(path: Path) -> ParsedDocument:
    import pptx

    presentation = pptx.Presentation(str(path))
    sections = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text.strip() for cell in row.cells))
        text = _nonempty("\n".join(parts))
        if text:
            sections.append(ParsedSection(text.splitlines()[0][:100], index, text, "slide"))
    if not sections:
        raise ValueError("PPTX 中没有可提取文本")
    return ParsedDocument(path.stem, "python-pptx", getattr(pptx, "__version__", "unknown"), sections)


def _parse_text(path: Path) -> ParsedDocument:
    raw = path.read_text(encoding="utf-8", errors="replace")
    blocks = [item.strip() for item in re.split(r"\n(?=#{1,6}\s|第[一二三四五六七八九十0-9]+[章节])", raw) if item.strip()]
    sections = []
    for index, block in enumerate(blocks or [raw], start=1):
        for offset in range(0, len(block), 3500):
            content = _nonempty(block[offset : offset + 3500])
            if content:
                title = content.splitlines()[0].lstrip("# ")[:100] or path.stem
                sections.append(ParsedSection(title, len(sections) + 1, content, "section"))
    if not sections:
        raise ValueError("文本文件为空")
    return ParsedDocument(path.stem, "utf-8-text", "1", sections)


def parse_document(path: Path) -> ParsedDocument:
    path = Path(path)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _parse_pdf(path)
    if suffix == ".docx":
        return _parse_docx(path)
    if suffix == ".pptx":
        return _parse_pptx(path)
    if suffix in {".txt", ".md", ".markdown"}:
        return _parse_text(path)
    raise ValueError(f"不支持的文件格式：{suffix or '无扩展名'}")


def render_index_text(document: ParsedDocument, document_id: str, course_id: str) -> str:
    rendered = []
    for section in document.sections:
        metadata = {
            "document_id": document_id,
            "course_id": course_id,
            "title": document.title,
            "chapter": section.chapter,
            "page_index": section.page_index,
            "source_kind": section.source_kind,
        }
        # Repeat metadata inside long sections so every embedding-sized chunk remains traceable.
        for offset in range(0, len(section.content), 320):
            fragment = section.content[offset : offset + 320].strip()
            if fragment:
                rendered.append(f"[{SOURCE_MARKER}]{json.dumps(metadata, ensure_ascii=False)}\n{fragment}")
    return "\n\n".join(rendered)


def extract_source_markers(content: str) -> List[dict]:
    pattern = re.compile(rf"\[{SOURCE_MARKER}\](\{{[^\n]+\}})")
    records = []
    for raw in pattern.findall(content or ""):
        try:
            value = json.loads(raw)
            if isinstance(value, dict):
                records.append(value)
        except json.JSONDecodeError:
            continue
    return records
